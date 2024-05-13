#!/usr/bin/env python
# -*- encoding: utf-8 -*-

import os
import pptx
import io
import csv
import pandas as pd

from PIL import Image
from typing import Iterable
from xagents.loader.utils import upload_to_minio, image2text
from xagents.loader.common import Chunk, AbstractLoader, ImageChunk, TextChunk, TableChunk, get_image_name_path, merge_chunks


def read_pptx_tab(tab, **kwargs):
    vf = io.StringIO()
    writer = csv.writer(vf)
    for row in tab.rows:
        writer.writerow(cell.text for cell in row.cells)
    vf.seek(0)
    return pd.read_csv(vf, **kwargs)


def analyse_slide(powerpoint:pptx.Presentation, file_name:str, upload_image:bool, ocr:bool)->Iterable[Chunk]:
    page_index = 1
    image_index = 0

    for slide in powerpoint.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:  # text
                content = shape.text_frame.text
                yield TextChunk(content=content, page_idx=page_index)
            
            elif shape.has_table:  # table
                df:pd.DataFrame = read_pptx_tab(shape.table)
                yield TableChunk(page_idx=page_index, data=df)
            
            elif shape.shape_type == 13:  # 图片类型，MSO_SHAPE_TYPE.PICTURE
                image_name, image_path = get_image_name_path(file_name, page_index, image_index)
                if upload_image:
                    with open(image_path, "wb") as f:
                        image = shape.image
                        image_bytes = image.blob
                        f.write(image_bytes)
                    url = upload_to_minio(image_path)
                else:
                    url=None
                content = image2text(Image.open(io.BytesIO(shape.image.blob))) if ocr else None
                image_index += 1
                yield ImageChunk(url=url, content=content,page_idx=page_index, image_name=image_name)
        page_index += 1


class PPTXLoader(AbstractLoader):
    def __init__(self, max_page:int=None, **kwargs):
        super().__init__(**kwargs)
        self.max_page = max_page
    
    def load(self, file_path: str, start_page=1, end_page=None, upload_image=True, ocr=False, **kwargs) -> Iterable[Chunk]:
        ppt = pptx.Presentation(file_path)
        chunks = analyse_slide(ppt,  file_name=os.path.basename(file_path), upload_image=upload_image, ocr=ocr)
        chunks = merge_chunks(chunks)
        yield from chunks


if __name__ == "__main__":
    ppt = pptx.Presentation("data/kb_file/时序数据预测.pptx")
    chunks = analyse_slide(ppt)
    for chunk in chunks:
        print(chunk)